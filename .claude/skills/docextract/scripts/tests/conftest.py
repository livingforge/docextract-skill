"""共有フィクスチャとテスト用ドキュメントのビルダー。

各形式 (docx/xlsx/pptx/pdf) を実ファイルとしてプログラム的に生成し、
抽出器を end-to-end で検証できるようにする。
"""

from __future__ import annotations

import zlib
from io import BytesIO
from pathlib import Path

import pytest


@pytest.fixture(scope="session")
def png_bytes() -> bytes:
    """小さな有効な PNG バイト列 (6x6, 赤)。

    ランタイム依存の Pillow で生成する (AGPL の PyMuPDF は使わない)。
    """
    from PIL import Image

    buf = BytesIO()
    Image.new("RGB", (6, 6), (200, 50, 50)).save(buf, format="PNG")
    return buf.getvalue()


@pytest.fixture
def png_file(tmp_path: Path, png_bytes: bytes) -> Path:
    p = tmp_path / "fixture.png"
    p.write_bytes(png_bytes)
    return p


# --------------------------------------------------------------------------
# docx ビルダー
# --------------------------------------------------------------------------
@pytest.fixture
def make_docx(tmp_path: Path):
    from docx import Document

    def _make(
        name: str = "doc.docx",
        *,
        paragraphs: list[tuple[str, str | None]] | None = None,
        table: list[list[str]] | None = None,
        image_path: Path | None = None,
        title: str | None = None,
        author: str | None = None,
    ) -> Path:
        d = Document()
        if title is not None:
            d.core_properties.title = title
        if author is not None:
            d.core_properties.author = author
        for text, style in paragraphs or []:
            if style:
                d.add_paragraph(text, style=style)
            else:
                d.add_paragraph(text)
        if table is not None:
            n_rows = len(table)
            n_cols = max((len(r) for r in table), default=0)
            t = d.add_table(rows=n_rows, cols=n_cols)
            for i, row in enumerate(table):
                for j, val in enumerate(row):
                    t.rows[i].cells[j].text = val
        if image_path is not None:
            d.add_picture(str(image_path))
        path = tmp_path / name
        path.parent.mkdir(parents=True, exist_ok=True)
        d.save(str(path))
        return path

    return _make


# --------------------------------------------------------------------------
# xlsx ビルダー
# --------------------------------------------------------------------------
def _inject_shapes(
    path: Path, shapes: dict[str, list[dict]], order: list[str]
) -> None:
    """保存済み xlsx にオートシェイプ/コネクタの drawing を直接注入する。

    openpyxl は autoshape を書き出せないため、zip を展開して
    xl/drawings/drawingK.xml と関連リレーションを手で組み立てる。
    各シェイプ dict: {"name","text","cell":(col,row)} = テキスト図形、
    {"connector": True} = コネクタ (テキスト無し、抽出対象外の確認用)。
    """
    import os
    import shutil
    import zipfile

    XDR = "http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing"
    AMAIN = "http://schemas.openxmlformats.org/drawingml/2006/main"
    RNS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

    def _sp(idx: int, s: dict) -> str:
        col, row = s.get("cell", (0, 0))
        to_col, to_row = s.get("to_cell", (col + 2, row + 2))
        anchor_from = (
            f"<xdr:from><xdr:col>{col}</xdr:col><xdr:colOff>0</xdr:colOff>"
            f"<xdr:row>{row}</xdr:row><xdr:rowOff>0</xdr:rowOff></xdr:from>"
            f"<xdr:to><xdr:col>{to_col}</xdr:col><xdr:colOff>0</xdr:colOff>"
            f"<xdr:row>{to_row}</xdr:row><xdr:rowOff>0</xdr:rowOff></xdr:to>"
        )
        if s.get("connector"):
            # st/end に接続先シェイプ id を渡すと明示接続 (<a:stCxn>/<a:endCxn>) になる
            cxn = ""
            if s.get("st") is not None:
                cxn += f'<a:stCxn id="{s["st"]}" idx="0"/>'
            if s.get("end") is not None:
                cxn += f'<a:endCxn id="{s["end"]}" idx="0"/>'
            body = (
                f'<xdr:cxnSp macro=""><xdr:nvCxnSpPr>'
                f'<xdr:cNvPr id="{idx}" name="{s.get("name","Conn")}"/>'
                f"<xdr:cNvCxnSpPr>{cxn}</xdr:cNvCxnSpPr></xdr:nvCxnSpPr>"
                f'<xdr:spPr><a:prstGeom prst="straightConnector1"><a:avLst/>'
                f"</a:prstGeom></xdr:spPr></xdr:cxnSp>"
            )
        else:
            paras = "".join(
                f"<a:p><a:r><a:t>{line}</a:t></a:r></a:p>"
                for line in str(s["text"]).split("\n")
            )
            body = (
                f'<xdr:sp macro="" textlink=""><xdr:nvSpPr>'
                f'<xdr:cNvPr id="{idx}" name="{s.get("name","Shape")}"/>'
                f"<xdr:cNvSpPr/></xdr:nvSpPr>"
                f'<xdr:spPr><a:prstGeom prst="rect"><a:avLst/></a:prstGeom></xdr:spPr>'
                f"<xdr:txBody><a:bodyPr/>{paras}</xdr:txBody></xdr:sp>"
            )
        return f"<xdr:twoCellAnchor>{anchor_from}{body}<xdr:clientData/></xdr:twoCellAnchor>"

    tmpd = path.parent / (path.stem + "_unzip")
    if tmpd.exists():
        shutil.rmtree(tmpd)
    with zipfile.ZipFile(path) as z:
        z.extractall(tmpd)

    ct_path = tmpd / "[Content_Types].xml"
    ct = ct_path.read_text(encoding="utf-8")
    sheet_index = {title: i + 1 for i, title in enumerate(order)}
    k = 0
    for sheet_name, sps in shapes.items():
        k += 1
        n = sheet_index[sheet_name]
        drawing = (
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            f'<xdr:wsDr xmlns:xdr="{XDR}" xmlns:a="{AMAIN}">'
            + "".join(_sp(i + 2, s) for i, s in enumerate(sps))
            + "</xdr:wsDr>"
        )
        (tmpd / "xl" / "drawings").mkdir(parents=True, exist_ok=True)
        (tmpd / "xl" / "drawings" / f"drawing{k}.xml").write_text(drawing, encoding="utf-8")

        sheet_path = tmpd / "xl" / "worksheets" / f"sheet{n}.xml"
        s = sheet_path.read_text(encoding="utf-8")
        if "xmlns:r=" not in s:
            s = s.replace("<worksheet ", f'<worksheet xmlns:r="{RNS}" ', 1)
        s = s.replace("</worksheet>", '<drawing r:id="rIdDraw"/></worksheet>')
        sheet_path.write_text(s, encoding="utf-8")

        rels_dir = tmpd / "xl" / "worksheets" / "_rels"
        rels_dir.mkdir(parents=True, exist_ok=True)
        rels_path = rels_dir / f"sheet{n}.xml.rels"
        rel = (
            f'<Relationship Id="rIdDraw" Type="{RNS}/drawing" '
            f'Target="../drawings/drawing{k}.xml"/>'
        )
        if rels_path.exists():
            r = rels_path.read_text(encoding="utf-8").replace(
                "</Relationships>", rel + "</Relationships>"
            )
        else:
            r = (
                '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
                '<Relationships xmlns="http://schemas.openxmlformats.org/'
                f'package/2006/relationships">{rel}</Relationships>'
            )
        rels_path.write_text(r, encoding="utf-8")

        override = (
            f'<Override PartName="/xl/drawings/drawing{k}.xml" '
            'ContentType="application/vnd.openxmlformats-officedocument.drawing+xml"/>'
        )
        if f"/xl/drawings/drawing{k}.xml" not in ct:
            ct = ct.replace("</Types>", override + "</Types>")
    ct_path.write_text(ct, encoding="utf-8")

    path.unlink()
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as z:
        for root, _, files in os.walk(tmpd):
            for f in files:
                fp = Path(root) / f
                z.write(fp, fp.relative_to(tmpd).as_posix())
    shutil.rmtree(tmpd)


@pytest.fixture
def make_xlsx(tmp_path: Path):
    from openpyxl import Workbook
    from openpyxl.drawing.image import Image as XLImage

    def _make(
        name: str = "book.xlsx",
        *,
        sheets: dict[str, list[list]] | None = None,
        merges: dict[str, list[str]] | None = None,  # sheet -> ["B1:B3", ...]
        image: tuple[str, Path, str] | None = None,  # (sheet, png_path, anchor)
        # sheet -> [{"name","text","cell":(col,row)} | {"connector":True}, ...]
        # openpyxl は autoshape を書けないため drawing XML を直接注入する。
        shapes: dict[str, list[dict]] | None = None,
        title: str | None = None,
        author: str | None = None,
    ) -> Path:
        wb = Workbook()
        # デフォルトシートを消す前に少なくとも 1 枚必要
        default = wb.active
        created_any = False
        for sheet_name, grid in (sheets or {}).items():
            if not created_any:
                ws = default
                ws.title = sheet_name
                created_any = True
            else:
                ws = wb.create_sheet(sheet_name)
            for r, row in enumerate(grid, start=1):
                for c, val in enumerate(row, start=1):
                    ws.cell(row=r, column=c, value=val)
        if not created_any:
            default.title = "Sheet1"
        for sheet_name, refs in (merges or {}).items():
            for ref in refs:
                wb[sheet_name].merge_cells(ref)
        if title is not None:
            wb.properties.title = title
        if author is not None:
            wb.properties.creator = author
        if image is not None:
            sheet_name, png_path, anchor = image
            ws = wb[sheet_name]
            ws.add_image(XLImage(str(png_path)), anchor)
        path = tmp_path / name
        path.parent.mkdir(parents=True, exist_ok=True)
        wb.save(str(path))
        if shapes:
            order = wb.sheetnames  # sheetN.xml は作成順 = この順
            _inject_shapes(path, shapes, order)
        return path

    return _make


# --------------------------------------------------------------------------
# pptx ビルダー
# --------------------------------------------------------------------------
@pytest.fixture
def make_pptx(tmp_path: Path):
    from pptx import Presentation
    from pptx.util import Inches

    def _make(
        name: str = "deck.pptx",
        *,
        slides: list[dict] | None = None,
        title: str | None = None,
        author: str | None = None,
    ) -> Path:
        prs = Presentation()
        blank = prs.slide_layouts[6]
        for spec in slides or []:
            slide = prs.slides.add_slide(blank)
            for text in spec.get("texts", []):
                tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
                tb.text_frame.text = text
            for tbl in spec.get("tables", []):
                n_rows = len(tbl)
                n_cols = max((len(r) for r in tbl), default=0)
                gt = slide.shapes.add_table(
                    n_rows, n_cols, Inches(1), Inches(3), Inches(5), Inches(2)
                ).table
                for i, row in enumerate(tbl):
                    for j, val in enumerate(row):
                        gt.cell(i, j).text = val
            for img in spec.get("images", []):
                slide.shapes.add_picture(str(img), Inches(5), Inches(1))
            notes = spec.get("notes")
            if notes is not None:
                slide.notes_slide.notes_text_frame.text = notes
        if title is not None:
            prs.core_properties.title = title
        if author is not None:
            prs.core_properties.author = author
        path = tmp_path / name
        path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(path))
        return path

    return _make


# --------------------------------------------------------------------------
# pdf ビルダー
#
# PyMuPDF (fitz, AGPL) を使わず、ランタイム依存の Pillow (画像デコード) と
# 標準ライブラリ zlib だけで最小の PDF をバイト列から直接組み立てる。
# 生成物は抽出器が読む要素を備える:
#   - テキスト層 (標準 Helvetica の Tj)           -> pdfplumber がテキスト行を抽出
#   - ストローク罫線 (m/l/S) + セル内テキスト      -> pdfplumber の罫線ベース表検出
#   - 埋め込み画像 XObject (FlateDecode/DeviceRGB) -> pypdf が画像として抽出
# 座標は旧 fitz フィクスチャに合わせ「左上原点・y 下向き」で受け取り、PDF の
# 「左下原点・y 上向き」へ変換する (Y = ページ高 - y)。
# --------------------------------------------------------------------------
_PDF_PAGE_W, _PDF_PAGE_H = 612.0, 792.0


def _pdf_num(v: float) -> bytes:
    """PDF 用の数値表記 (整数はそのまま、少数は末尾の 0 を落とす)。"""
    return f"{v:.2f}".rstrip("0").rstrip(".").encode("latin-1")


def _pdf_str(s: str) -> bytes:
    """PDF リテラル文字列。特殊文字 \\ ( ) をエスケープする。"""
    out = s.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    return out.encode("latin-1", "replace")


def _pdf_image_xobject(img_bytes: bytes) -> bytes:
    """画像バイト列を FlateDecode/DeviceRGB の画像 XObject 本体に変換する。"""
    from PIL import Image

    img = Image.open(BytesIO(img_bytes)).convert("RGB")
    w, h = img.size
    comp = zlib.compress(img.tobytes())
    return (
        b"<</Type/XObject/Subtype/Image/Width %d/Height %d"
        b"/ColorSpace/DeviceRGB/BitsPerComponent 8/Filter/FlateDecode"
        b"/Length %d>>\nstream\n%s\nendstream" % (w, h, len(comp), comp)
    )


def _pdf_page_content(spec: dict) -> bytes:
    """1 ページ分のコンテンツストリーム (テキスト・罫線・画像描画命令) を作る。"""
    H = _PDF_PAGE_H
    parts: list[bytes] = []
    for text, (x, y) in spec.get("texts", []):
        parts.append(
            b"BT /F1 11 Tf %s %s Td (%s) Tj ET\n"
            % (_pdf_num(x), _pdf_num(H - y), _pdf_str(text))
        )
    grid = spec.get("grid")
    if grid:
        n_rows, n_cols = grid["rows"], grid["cols"]
        x0, y0 = grid.get("origin", (100, 200))
        cw, ch = grid.get("cw", 80), grid.get("ch", 30)
        for i in range(n_rows + 1):  # 水平罫線
            yy = H - (y0 + i * ch)
            parts.append(b"%s %s m %s %s l S\n" % (
                _pdf_num(x0), _pdf_num(yy),
                _pdf_num(x0 + n_cols * cw), _pdf_num(yy)))
        for j in range(n_cols + 1):  # 垂直罫線
            xx = x0 + j * cw
            parts.append(b"%s %s m %s %s l S\n" % (
                _pdf_num(xx), _pdf_num(H - y0),
                _pdf_num(xx), _pdf_num(H - (y0 + n_rows * ch))))
        for (i, j), val in grid.get("cells", {}).items():
            parts.append(
                b"BT /F1 10 Tf %s %s Td (%s) Tj ET\n"
                % (_pdf_num(x0 + j * cw + 4), _pdf_num(H - (y0 + i * ch + 18)),
                   _pdf_str(val)))
    for idx, (_img, rect) in enumerate(spec.get("images", []), start=1):
        x0, y0, x1, y1 = rect
        parts.append(b"q %s 0 0 %s %s %s cm /Im%d Do Q\n" % (
            _pdf_num(x1 - x0), _pdf_num(y1 - y0),
            _pdf_num(x0), _pdf_num(H - y1), idx))
    return b"".join(parts)


def _write_pdf(path: Path, pages: list[dict] | None,
               title: str | None, author: str | None) -> None:
    """最小構成の PDF をバイト列から組み立てて path に書き出す。"""
    pages = pages if pages else [{}]

    # オブジェクト番号を先に割り当てる (Pages の Kids が前方参照になるため)。
    # 1=Catalog, 2=Pages, 3=Font, 4=Info、以降ページごとに content/画像/page。
    body: dict[int, bytes] = {
        1: b"<</Type/Catalog/Pages 2 0 R>>",
        3: b"<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>",
    }
    info = b"<<"
    if title is not None:
        info += b"/Title(%s)" % _pdf_str(title)
    if author is not None:
        info += b"/Author(%s)" % _pdf_str(author)
    body[4] = info + b">>"

    next_obj = 5
    page_nums: list[int] = []
    for spec in pages:
        content = _pdf_page_content(spec)
        content_num = next_obj
        next_obj += 1
        body[content_num] = (
            b"<</Length %d>>\nstream\n%s\nendstream" % (len(content), content)
        )
        image_nums: list[int] = []
        for img_bytes, _rect in spec.get("images", []):
            body[next_obj] = _pdf_image_xobject(img_bytes)
            image_nums.append(next_obj)
            next_obj += 1
        res = b"/Font<</F1 3 0 R>>"
        if image_nums:
            xobjs = b"".join(
                b"/Im%d %d 0 R" % (k, num)
                for k, num in enumerate(image_nums, start=1))
            res += b"/XObject<<%s>>" % xobjs
        page_num = next_obj
        next_obj += 1
        body[page_num] = (
            b"<</Type/Page/Parent 2 0 R/MediaBox[0 0 %s %s]"
            b"/Resources<<%s>>/Contents %d 0 R>>" % (
                _pdf_num(_PDF_PAGE_W), _pdf_num(_PDF_PAGE_H),
                res, content_num))
        page_nums.append(page_num)

    kids = b" ".join(b"%d 0 R" % n for n in page_nums)
    body[2] = b"<</Type/Pages/Kids[%s]/Count %d>>" % (kids, len(page_nums))

    # 本体を番号順に直列化し、各オブジェクトの先頭オフセットを記録する。
    buf = bytearray(b"%PDF-1.4\n%\xe2\xe3\xcf\xd3\n")
    offsets: dict[int, int] = {}
    for num in range(1, next_obj):
        offsets[num] = len(buf)
        buf += b"%d 0 obj\n" % num
        buf += body[num]
        buf += b"\nendobj\n"

    xref_off = len(buf)
    size = next_obj  # オブジェクト 0..next_obj-1
    buf += b"xref\n0 %d\n" % size
    buf += b"0000000000 65535 f \n"
    for num in range(1, size):
        buf += b"%010d 00000 n \n" % offsets[num]
    buf += b"trailer\n<</Size %d/Root 1 0 R/Info 4 0 R>>\n" % size
    buf += b"startxref\n%d\n%%%%EOF\n" % xref_off

    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_bytes(bytes(buf))


@pytest.fixture
def make_pdf(tmp_path: Path):
    def _make(
        name: str = "doc.pdf",
        *,
        pages: list[dict] | None = None,
        title: str | None = None,
        author: str | None = None,
    ) -> Path:
        path = tmp_path / name
        _write_pdf(path, pages, title, author)
        return path

    return _make
