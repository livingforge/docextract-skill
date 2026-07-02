"""docextract の評価ランナー (eval runner)。

同梱テスト (``scripts/tests/``) が「自コードの自己検証」なのに対し、こちらは
**評価データセットと合否基準を data として外出し** (``cases.jsonl``) し、それを
列挙実行して合否を集計する視点分離の評価ハーネスである。基準 (期待値・しきい値・
必須フィールド) はコードに埋め込まず ``cases.jsonl`` の ``expect`` ブロックで宣言する。

使い方::

    python run_eval.py                 # 同ディレクトリの cases.jsonl を実行
    python run_eval.py path/to/cases.jsonl
    python run_eval.py --json          # 集計を機械可読な JSON で出力

各ケースは入力フィクスチャを **実行時に生成** (docx/xlsx/pptx を runtime 依存だけで
構築) するため、ネットワークも OCR モデルも不要で決定論的。抽出は ``ocr=False`` /
``image_tables=False`` (非決定なパスを無効化) で走らせ、宣言した合否基準と突き合わせる。

合否基準 (``expect`` に data として宣言できるもの)::

    required_top_keys : result.json に必須のトップレベルキー
    summary_min       : 要素種別ごとの最小件数 (例 {"text": 3, "table": 1})
    must_contain_text : いずれかの text/table 要素に含まれるべき文字列
    no_degradations   : true なら劣化痕跡ゼロを要求 (取りこぼしなし)

終了コード: 全ケース pass なら 0、1 件でも fail なら 1。
"""

from __future__ import annotations

import argparse
import json
import sys
import tempfile
from pathlib import Path


def _bootstrap_docextract() -> None:
    """docextract パッケージのある場所を探して sys.path に載せる。

    - 配布バンドル: ``scripts/eval/`` の隣 ``scripts/docextract`` (parents[1])。
    - リポジトリ: ルート直下の ``docextract`` (数階層上)。
    両レイアウトで動くよう、スクリプトから上へ ``docextract/__init__.py`` を探す。
    """
    here = Path(__file__).resolve()
    for base in [here.parent, *here.parents]:
        if (base / "docextract" / "__init__.py").is_file():
            sys.path.insert(0, str(base))
            return
    # 見つからなければ parents[1] を入れて後段の import エラーに委ねる。
    sys.path.insert(0, str(here.parents[1]))


_bootstrap_docextract()

from docextract import extract  # noqa: E402


# ── フィクスチャ生成 (runtime 依存だけで構築、PDF は PyMuPDF 非依存のため対象外) ──
def _build_docx(spec: dict, path: Path) -> None:
    from docx import Document

    d = Document()
    for text, style in spec.get("paragraphs", []) or []:
        d.add_paragraph(text, style=style) if style else d.add_paragraph(text)
    table = spec.get("table")
    if table:
        n_rows, n_cols = len(table), max((len(r) for r in table), default=0)
        t = d.add_table(rows=n_rows, cols=n_cols)
        for i, row in enumerate(table):
            for j, val in enumerate(row):
                t.rows[i].cells[j].text = str(val)
    d.save(str(path))


def _build_xlsx(spec: dict, path: Path) -> None:
    from openpyxl import Workbook

    wb = Workbook()
    default = wb.active
    created = False
    for name, grid in (spec.get("sheets") or {}).items():
        ws = default if not created else wb.create_sheet(name)
        if not created:
            ws.title = name
            created = True
        for r, row in enumerate(grid, start=1):
            for c, val in enumerate(row, start=1):
                ws.cell(row=r, column=c, value=val)
    wb.save(str(path))


def _build_pptx(spec: dict, path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[6]
    for slide_spec in spec.get("slides", []) or []:
        slide = prs.slides.add_slide(blank)
        for text in slide_spec.get("texts", []):
            tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
            tb.text_frame.text = text
        for tbl in slide_spec.get("tables", []):
            n_rows, n_cols = len(tbl), max((len(r) for r in tbl), default=0)
            gt = slide.shapes.add_table(
                n_rows, n_cols, Inches(1), Inches(3), Inches(5), Inches(2)
            ).table
            for i, row in enumerate(tbl):
                for j, val in enumerate(row):
                    gt.cell(i, j).text = str(val)
    prs.save(str(path))


_BUILDERS = {"docx": _build_docx, "xlsx": _build_xlsx, "pptx": _build_pptx}


# ── 合否判定 (基準は case の expect に data として宣言される) ──
def _element_texts(data: dict) -> list[str]:
    out: list[str] = []
    for el in data.get("elements", []):
        if el.get("type") == "text":
            out.append(el.get("content", ""))
        elif el.get("type") == "table":
            out.extend(str(c) for row in el.get("rows", []) for c in row)
    return out


def _check(data: dict, expect: dict) -> list[str]:
    """宣言された基準に対する違反理由のリストを返す (空なら pass)。"""
    failures: list[str] = []

    for key in expect.get("required_top_keys", []):
        if key not in data:
            failures.append(f"必須キー欠落: {key}")

    summary = data.get("summary", {})
    for kind, minimum in (expect.get("summary_min") or {}).items():
        got = summary.get(kind, 0)
        if got < minimum:
            failures.append(f"{kind} 件数不足: {got} < {minimum}")

    haystack = "\n".join(_element_texts(data))
    for needle in expect.get("must_contain_text", []):
        if needle not in haystack:
            failures.append(f"本文に見当たらない: {needle!r}")

    if expect.get("no_degradations"):
        deg = data.get("degraded", {}).get("count", 0)
        if deg:
            failures.append(f"劣化痕跡あり: {deg} 件")

    return failures


def run_case(case: dict, workdir: Path) -> dict:
    fmt = case["format"]
    builder = _BUILDERS.get(fmt)
    if builder is None:
        return {"id": case["id"], "passed": False, "failures": [f"未対応の format: {fmt}"]}
    src = workdir / f"{case['id']}.{fmt}"
    builder(case.get("build", {}), src)
    opts = case.get("extract", {})
    data = extract(
        src,
        output_dir=workdir / "out",
        ocr=opts.get("ocr", False),
        image_tables=opts.get("image_tables", False),
        record_manifest=False,
    )
    failures = _check(data, case.get("expect", {}))
    return {"id": case["id"], "passed": not failures, "failures": failures}


def load_cases(path: Path) -> list[dict]:
    cases = []
    for line in path.read_text(encoding="utf-8").splitlines():
        line = line.strip()
        if line and not line.startswith("#"):
            cases.append(json.loads(line))
    return cases


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="docextract 評価ランナー")
    parser.add_argument(
        "cases",
        nargs="?",
        default=str(Path(__file__).resolve().parent / "cases.jsonl"),
        help="評価ケース (JSON Lines)。既定は同ディレクトリの cases.jsonl",
    )
    parser.add_argument("--json", action="store_true", help="集計を JSON で出力")
    args = parser.parse_args(argv)

    for stream in (sys.stdout, sys.stderr):
        try:
            stream.reconfigure(encoding="utf-8")
        except (AttributeError, ValueError):
            pass

    cases = load_cases(Path(args.cases))
    results = []
    with tempfile.TemporaryDirectory() as tmp:
        workdir = Path(tmp)
        for case in cases:
            try:
                results.append(run_case(case, workdir))
            except Exception as e:  # noqa: BLE001 - ケース単位の異常を fail として集計
                results.append(
                    {"id": case.get("id", "?"), "passed": False, "failures": [f"例外: {e!r}"]}
                )

    passed = sum(1 for r in results if r["passed"])
    summary = {"total": len(results), "passed": passed, "failed": len(results) - passed, "cases": results}

    if args.json:
        print(json.dumps(summary, ensure_ascii=False, indent=2))
    else:
        for r in results:
            mark = "PASS" if r["passed"] else "FAIL"
            print(f"[{mark}] {r['id']}")
            for f in r["failures"]:
                print(f"        - {f}")
        print(f"\n合計 {summary['total']} / 成功 {passed} / 失敗 {summary['failed']}")

    return 0 if summary["failed"] == 0 else 1


if __name__ == "__main__":
    raise SystemExit(main())
